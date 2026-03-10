from pptx import Presentation
from pptx.util import Inches, Pt
import os

def generate_pptx(slides_file="data/slides_outline.txt", images_dir="output/images", out_path="output/presentation.pptx"):
    prs = Presentation()

    with open(slides_file, "r", encoding="utf-8") as f:
        slides = f.read().split("Slide")

    for i, slide_block in enumerate(slides, start=1):
        if slide_block.strip():
            lines = [line.strip() for line in slide_block.split("\n") if line.strip()]

            title = ""
            subtitle = ""
            points = []
            for line in lines:
                if line.startswith("Title:"):
                    title = line.replace("Title:", "").strip()
                elif line.startswith("Subtitle:"):
                    subtitle = line.replace("Subtitle:", "").strip()
                elif line.startswith("-"):
                    points.append(line)

            slide = prs.slides.add_slide(prs.slide_layouts[5])  
            slide.shapes.title.text = title

            if subtitle:
                txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(1))
                tf = txBox.text_frame
                p = tf.add_paragraph()
                p.text = subtitle
                p.font.size = Pt(16)
                p.font.italic = True
                p.font.color.rgb = (100,100,100)

            body_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(4.5), Inches(3))
            tf = body_box.text_frame
            for pt in points:
                p = tf.add_paragraph()
                p.text = pt.replace("-", "").strip()
                p.level = 0

            img_path = os.path.join(images_dir, f"slide{i}.png")
            if os.path.exists(img_path):
                left = Inches(5)
                top = Inches(1.5)
                slide.shapes.add_picture(img_path, left, top, width=Inches(4.5))

    prs.save(out_path)
    print(f"[INFO] Presentation saved to {out_path}")