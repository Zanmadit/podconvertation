import os
import shutil
import subprocess
import glob
import sys
from pathlib import Path

from gtts import gTTS
from moviepy.editor import ImageClip, AudioFileClip, concatenate_videoclips
from pptx import Presentation

def check_binary(name):
    from shutil import which
    return which(name) is not None

def pptx_to_pdf(pptx_path: str, outdir: str):
    if not check_binary("libreoffice") and not check_binary("soffice"):
        raise RuntimeError("LibreOffice (soffice) not found on PATH. Install libreoffice.")
    bin_name = "soffice" if check_binary("soffice") else "libreoffice"
    cmd = [
        bin_name,
        "--headless",
        "--convert-to",
        "pdf",
        "--outdir",
        outdir,
        pptx_path,
    ]
    print("[INFO] Running:", " ".join(cmd))
    subprocess.run(cmd, check=True)
    pdf_name = Path(pptx_path).with_suffix(".pdf").name
    pdf_path = os.path.join(outdir, pdf_name)
    if not os.path.exists(pdf_path):
        raise RuntimeError(f"LibreOffice failed to produce PDF at {pdf_path}")
    return pdf_path

def pdf_to_pngs(pdf_path: str, outdir: str, prefix="slide", dpi=150):
    """Convert PDF -> multiple PNGs using pdftoppm (poppler)"""
    if not check_binary("pdftoppm"):
        raise RuntimeError("pdftoppm not found on PATH. Install poppler (pdftoppm).")
    out_prefix = os.path.join(outdir, prefix)
    cmd = ["pdftoppm", "-png", "-r", str(dpi), pdf_path, out_prefix]
    print("[INFO] Running:", " ".join(cmd))
    subprocess.run(cmd, check=True)
    pngs = sorted(glob.glob(f"{out_prefix}-*.png"))
    return pngs

def extract_slide_texts(pptx_path: str):
    prs = Presentation(pptx_path)
    slides_texts = []
    for slide in prs.slides:
        title_text = ""
        try:
            if slide.shapes.title and slide.shapes.title.text:
                title_text = slide.shapes.title.text.strip()
        except Exception:
            title_text = ""
        bullets = []
        for shape in slide.shapes:
            try:
                if shape == slide.shapes.title:
                    continue
            except Exception:
                pass
            if not getattr(shape, "has_text_frame", False):
                continue
            try:
                txt = shape.text.strip()
                if txt:
                    for para in txt.splitlines():
                        p = para.strip()
                        if p:
                            bullets.append(p)
            except Exception:
                continue
        pieces = []
        if title_text:
            pieces.append(title_text)
        if bullets:
            pieces.append(" ".join(bullets[:3]))
        narration = ". ".join(pieces).strip()
        if not narration:
            narration = "Slide."
        slides_texts.append(narration)
    return slides_texts

def make_video_from_pptx(
    pptx_file="output/presentation.pptx",
    out_dir="output",
    lang="en",
    cleanup=True,
    dpi=150,
):
    """
    Full flow:
     - pptx -> pdf (libreoffice)
     - pdf -> pngs (pdftoppm)
     - extract texts from pptx -> per-slide narration
     - tts (gTTS) -> mp3 per slide
     - moviepy: image+audio -> final mp4
    """
    out_dir = os.path.abspath(out_dir)
    os.makedirs(out_dir, exist_ok=True)

    slides_png_dir = os.path.join(out_dir, "slides_png")
    audio_dir = os.path.join(out_dir, "audio")
    os.makedirs(slides_png_dir, exist_ok=True)
    os.makedirs(audio_dir, exist_ok=True)

    pptx_path = os.path.abspath(pptx_file)
    if not os.path.exists(pptx_path):
        raise FileNotFoundError(f"PPTX not found: {pptx_path}")

    print("[STEP] Converting PPTX -> PDF")
    pdf_path = pptx_to_pdf(pptx_path, outdir=slides_png_dir)

    print("[STEP] Converting PDF -> PNGs")
    png_list = pdf_to_pngs(pdf_path, slides_png_dir, prefix="slide", dpi=dpi)
    if not png_list:
        raise RuntimeError("No PNGs were created from PDF")

    print("[STEP] Extracting slide texts from PPTX")
    texts = extract_slide_texts(pptx_path)
    print(f"[INFO] Found {len(texts)} slides in PPTX, {len(png_list)} PNG images exported")

    n = min(len(texts), len(png_list))
    if n == 0:
        raise RuntimeError("No slides/images available to create video")

    clips = []
    for i in range(n):
        img = png_list[i]
        narration = texts[i]
        audio_path = os.path.join(audio_dir, f"slide{i+1}.mp3")
        print(f"[INFO] Generating TTS for slide {i+1}: {narration[:120]}")

        tts = gTTS(text=narration, lang=lang)
        tts.save(audio_path)

        audio_clip = AudioFileClip(audio_path)
        clip = (
            ImageClip(img)
            .set_duration(audio_clip.duration + 0.2)  
            .set_audio(audio_clip)
            .fadein(0.35)
            .fadeout(0.35)
        )
        clips.append(clip)

    print("[STEP] Concatenating clips into final video")
    final = concatenate_videoclips(clips, method="compose")
    out_video = os.path.join(out_dir, "presentation.mp4")
    final.write_videofile(out_video, fps=24, codec="libx264", audio_codec="aac")
    print(f"[DONE] Video saved to {out_video}")

    if cleanup:
        try:
            if os.path.exists(pdf_path):
                os.remove(pdf_path)

        except Exception as e:
            print("[WARN] cleanup failed:", e)

    return out_video

if __name__ == "__main__":
    try:
        make_video_from_pptx(pptx_file="output/presentation.pptx", out_dir="output", lang="en", cleanup=False)
    except Exception as e:
        print("[ERROR]", e)
        sys.exit(1)
